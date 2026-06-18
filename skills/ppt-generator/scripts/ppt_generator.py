#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT 生成器 — 两阶段AI生成 + python-pptx 构建
Stage 1: 火山引擎豆包模型生成 PPT 结构大纲
Stage 2: 逐页丰富内容（详细要点、关键数据、演讲备注）
Stage 3: python-pptx 渲染高质量 PPT
"""

import sys
import os
import json
import subprocess

# ============================================================
# 依赖
# ============================================================

def install_package(pkg):
    try:
        subprocess.check_call([
            sys.executable, "-m", "pip", "install", pkg, "--break-system-packages"
        ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        return True
    except Exception:
        return False

def check_dependencies():
    for pkg in ["requests", "python-pptx"]:
        mod = pkg.replace("-", "_")
        try:
            __import__(mod)
        except ImportError:
            if not install_package(pkg):
                print("Error: 无法安装 " + pkg)
                return False
    return True

# ============================================================
# 火山引擎 API 调用
# ============================================================

def get_api_config():
    base_url = os.environ.get("VOLCENGINE_BASE_URL", "")
    api_key  = os.environ.get("VOLCENGINE_API_KEY", "")
    model    = os.environ.get("VOLCENGINE_MODEL", "")
    if not base_url or not api_key:
        return None
    return {
        "url": base_url.rstrip("/") + "/chat/completions",
        "api_key": api_key,
        "model": model
    }

def call_llm(config, messages, max_tokens=8000, temperature=0.6):
    import requests
    headers = {
        "Content-Type": "application/json",
        "Authorization": "Bearer " + config["api_key"]
    }
    body = {
        "model": config["model"],
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens
    }
    resp = requests.post(config["url"], headers=headers, json=body, timeout=120)
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]

def extract_json(text):
    """从AI响应中提取JSON，兼容被markdown包裹的情况"""
    text = text.strip()
    # 尝试 ```json ... ```
    for prefix in ["```json", "```JSON", "```"]:
        start = text.find(prefix)
        if start != -1:
            start += len(prefix)
            end = text.find("```", start)
            if end != -1:
                return json.loads(text[start:end].strip())
    # 尝试直接解析
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    # 尝试找到最外层 { }
    brace_count = 0
    start = -1
    for i, ch in enumerate(text):
        if ch == '{':
            if brace_count == 0:
                start = i
            brace_count += 1
        elif ch == '}':
            brace_count -= 1
            if brace_count == 0 and start != -1:
                try:
                    return json.loads(text[start:i+1])
                except json.JSONDecodeError:
                    start = -1
    raise ValueError("无法从AI响应中提取JSON")

# ============================================================
# Stage 1: 生成结构大纲
# ============================================================

SYSTEM_PROMPT_STRUCTURE = """你是一位资深的PPT架构师和内容策划专家，擅长将复杂信息转化为清晰、有逻辑、有说服力的演示文稿。

你的任务是分析用户提供的文本内容，设计一份高质量的PPT结构大纲。

## 设计原则

1. **逻辑为王**：页面之间必须有清晰的逻辑递进关系，不能是简单的信息罗列
2. **黄金结构**：
   - 开篇：封面 → 目录/议程 → 背景/现状（为什么要做这件事）
   - 主体：核心内容分层展开，每页聚焦一个主题
   - 收尾：总结要点 → 下一步/行动计划 → 致谢
3. **每页聚焦**：每页只讲一个核心观点，用 3-5 个要点支撑
4. **要点精炼**：每个要点控制在 15-30 字，用动词开头，避免流水账
5. **数据驱动**：如果原文有数据、指标、对比，必须提炼进要点
6. **标题有力**：每页标题是核心结论而非主题描述（如用"用户增长3倍"而非"用户增长情况"）

## 页数指南
- 简单内容（500字以下）：5-8 页
- 中等内容（500-2000字）：8-12 页
- 复杂内容（2000字以上）：12-18 页

## 输出要求
- 只输出JSON，不要任何解释
- JSON格式严格如下"""

USER_PROMPT_STRUCTURE = """请为以下内容设计PPT结构大纲：

---
{content}
---

输出严格JSON格式：
{{
  "title": "PPT主标题 — 简洁有力的核心主张",
  "subtitle": "副标题 — 补充说明或日期",
  "theme_color": "深蓝色",
  "slides": [
    {{
      "page": 1,
      "type": "cover",
      "header": "主标题",
      "subheader": "副标题或演讲者信息"
    }},
    {{
      "page": 2,
      "type": "agenda",
      "header": "今日议程",
      "items": ["一、背景与现状", "二、核心方案", "三、实施计划"]
    }},
    {{
      "page": 3,
      "type": "content",
      "header": "核心观点（用结论做标题）",
      "subheader": "一句话补充说明",
      "bullets": [
        "要点一：动词开头，15-30字，包含关键数据",
        "要点二：...",
        "要点三：..."
      ]
    }}
  ]
}}

type 可选值：
- cover：封面页
- agenda：目录/议程页
- content：常规内容页
- data：数据/对比页（bullets 中放数据点）
- summary：总结页
- action：行动计划/下一步
- thanks：致谢页

要求：
1. 最后必须有 summary 或 action 页 + thanks 页
2. 内容页的 header 必须是结论性的，不要用"介绍""概述"等平庸标题
3. 根据内容长度自动决定页数"""

# ============================================================
# Stage 2: 逐页丰富内容
# ============================================================

SYSTEM_PROMPT_ENRICH = """你是一位专业的PPT内容编辑。你的任务是丰富单页PPT的内容，使其信息密度更高、表达更精炼、更有说服力。"""

USER_PROMPT_ENRICH = """请丰富以下PPT页面的内容：

原始页面信息：
{slide_json}

页面上下文（前后页内容概要）：
{context}

请输出丰富后的JSON，格式保持一致，但要做到：
1. bullets 扩展到 4-6 个，每个要点更有信息量
2. 添加 "key_data" 字段：列出本页涉及的关键数据或数字（如 "增长200%"、"覆盖50万用户"）
3. 添加 "speaker_notes" 字段：50-100字的演讲备注，说明这页该怎么讲
4. 如果是 data 类型，添加 "comparison" 字段：对比维度的描述

只输出JSON，不要任何解释。"""

# ============================================================
# Stage 1+2: 生成完整大纲
# ============================================================

def generate_outline(config, user_content):
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT_STRUCTURE},
        {"role": "user", "content": USER_PROMPT_STRUCTURE.format(content=user_content)}
    ]

    print("[1/3] 正在生成PPT结构大纲...")
    raw = call_llm(config, messages, max_tokens=8000, temperature=0.5)
    outline = extract_json(raw)

    slides = outline.get("slides", [])
    total = len(slides)
    print("      大纲生成完成，共 " + str(total) + " 页")

    # Stage 2: 逐页丰富（跳过 cover / thanks）
    enrichable = [s for s in slides if s.get("type") not in ("cover", "thanks")]
    if enrichable:
        print("[2/3] 正在丰富 " + str(len(enrichable)) + " 页内容...")
        for i, slide in enumerate(slides):
            if slide.get("type") in ("cover", "thanks"):
                continue
            try:
                # 构建上下文
                idx = slides.index(slide)
                prev_header = slides[idx-1]["header"] if idx > 0 else "（首页）"
                next_header = slides[idx+1]["header"] if idx < len(slides)-1 else "（末页）"
                context = "上一页: " + prev_header + "\n下一页: " + next_header

                msg = [
                    {"role": "system", "content": SYSTEM_PROMPT_ENRICH},
                    {"role": "user", "content": USER_PROMPT_ENRICH.format(
                        slide_json=json.dumps(slide, ensure_ascii=False, indent=2),
                        context=context
                    )}
                ]
                enriched_raw = call_llm(config, msg, max_tokens=4000, temperature=0.4)
                enriched = extract_json(enriched_raw)
                slides[i] = enriched
                print("      第" + str(i+1) + "页已丰富")
            except Exception as e:
                print("      第" + str(i+1) + "页丰富失败，保留原始: " + str(e))
    else:
        print("[2/3] 无需丰富（仅封面/致谢页）")

    outline["slides"] = slides
    return outline

# ============================================================
# Stage 3: python-pptx 构建
# ============================================================

def build_pptx(outline, output_path):
    from pptx import Presentation
    from pptx.util import Pt, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn

    print("[3/3] 正在构建PPT文件...")

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 配色
    COLOR_DARK   = RGBColor(0x1A, 0x25, 0x3C)   # 深蓝黑
    COLOR_ACCENT = RGBColor(0x21, 0x73, 0x46)   # 绿色强调
    COLOR_GRAY   = RGBColor(0x66, 0x66, 0x66)
    COLOR_LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
    COLOR_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

    def set_font(run, size=18, bold=False, color=COLOR_DARK):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Microsoft YaHei"

    def add_shape_bg(slide, color):
        """给幻灯片加背景色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    title_text = outline.get("title", "演示文稿")
    subtitle_text = outline.get("subtitle", "")
    slides_data = outline.get("slides", [])

    for s in slides_data:
        stype = s.get("type", "content")

        # ---- 封面 ----
        if stype == "cover":
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            add_shape_bg(slide, COLOR_DARK)

            # 标题
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s.get("header", title_text)
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

            # 副标题
            if s.get("subheader", subtitle_text):
                p2 = tf.add_paragraph()
                p2.text = s.get("subheader", subtitle_text)
                p2.font.size = Pt(20)
                p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(24)

        # ---- 议程/目录 ----
        elif stype == "agenda":
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # 左侧色块
            slide.shapes.add_shape(
                1, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
            ).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = COLOR_ACCENT

            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(11), Inches(1))
            p = txBox.text_frame.paragraphs[0]
            p.text = s.get("header", "议程")
            set_font(p.add_run(), 36, True, COLOR_DARK)
            p.alignment = PP_ALIGN.LEFT

            items = s.get("items", s.get("bullets", []))
            for j, item in enumerate(items):
                txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(2.2 + j * 1.1), Inches(10), Inches(0.8))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                set_font(p2.add_run(), 24, False, COLOR_DARK)
                p2.text = item

        # ---- 总结/行动/致谢 ----
        elif stype in ("summary", "action", "thanks"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_shape_bg(slide, COLOR_DARK)

            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = s.get("header", "谢谢")
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

            bullets = s.get("bullets", [])
            if bullets:
                txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.2), Inches(10), Inches(3.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for j, b in enumerate(bullets):
                    p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    p2.text = "  " + b
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
                    p2.space_after = Pt(10)
                    p2.alignment = PP_ALIGN.LEFT

            notes = s.get("speaker_notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        # ---- 内容页 / 数据页 ----
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # 顶部色条
            bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLOR_ACCENT

            # 标题
            header = s.get("header", "")
            subheader = s.get("subheader", "")

            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = header
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = COLOR_DARK
            p.alignment = PP_ALIGN.LEFT

            if subheader:
                p2 = tf.add_paragraph()
                p2.text = subheader
                p2.font.size = Pt(16)
                p2.font.color.rgb = COLOR_GRAY
                p2.space_before = Pt(4)

            # 要点列表
            bullets = s.get("bullets", [])
            if bullets:
                txBox2 = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11), Inches(4.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for j, b in enumerate(bullets):
                    p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    run = p2.add_run()
                    run.text = "●  " + b
                    run.font.size = Pt(18)
                    run.font.color.rgb = COLOR_DARK
                    p2.space_after = Pt(10)
                    p2.level = 0

            # 关键数据高亮
            key_data = s.get("key_data", [])
            if key_data:
                y_pos = min(1.8 + len(bullets) * 0.65 + 0.3, 5.5)
                txBox3 = slide.shapes.add_textbox(Inches(1.0), Inches(y_pos), Inches(11), Inches(1.2))
                tf3 = txBox3.text_frame
                tf3.word_wrap = True
                p3 = tf3.paragraphs[0]
                p3.text = "关键数据: " + "  |  ".join(key_data) if isinstance(key_data, list) else str(key_data)
                p3.font.size = Pt(14)
                p3.font.color.rgb = COLOR_ACCENT
                p3.font.bold = True

            # 演讲备注
            notes = s.get("speaker_notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_path)
    return output_path

# ============================================================
# Main
# ============================================================

def main():
    args = sys.argv[1:]

    if args and args[0] == "test":
        print("test ok")
        return

    if not args:
        print("Error: 请提供生成PPT的内容")
        return

    content = " ".join(args)

    if not check_dependencies():
        return

    config = get_api_config()
    if config is None:
        print("Error: 未检测到火山引擎配置（VOLCENGINE_BASE_URL / VOLCENGINE_API_KEY）")
        return

    print("使用模型: " + config["model"])
    print("输入内容: " + str(len(content)) + " 字符")
    print()

    try:
        outline = generate_outline(config, content)
    except Exception as e:
        print("Error: AI生成大纲失败: " + str(e))
        return

    output_path = os.path.join(os.getcwd(), "generated_presentation.pptx")

    try:
        build_pptx(outline, output_path)
    except Exception as e:
        print("Error: 构建PPT失败: " + str(e))
        return

    slides = outline.get("slides", [])
    print()
    print("=" * 50)
    print("PPT生成成功！")
    print("文件: " + output_path)
    print("标题: " + outline.get("title", ""))
    subtitle = outline.get("subtitle", "")
    if subtitle:
        print("副标题: " + subtitle)
    print("页数: " + str(len(slides)))
    print("页类型分布: " + ", ".join(
        [s.get("type", "content") for s in slides]
    ))
    print("=" * 50)

if __name__ == "__main__":
    main()
